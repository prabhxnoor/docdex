"""Deterministic synthetic corpus for the docdex value benchmark.

Models the real-world problem docdex exists for: facts buried inside Office
documents and PDFs whose filenames lie about their contents, surrounded by
distractor files that share vocabulary with the queries.

Everything is seeded — two runs produce byte-identical corpora, so results
are reproducible across machines.
"""
from __future__ import annotations

import random
import zlib
from pathlib import Path

SEED = 42

WORDS = (
    "agreement vendor invoice quarterly review board approval compliance "
    "budget travel policy renewal liability procurement onboarding security "
    "audit milestone deliverable retention payment schedule warranty notice "
    "termination escalation committee headcount forecast pipeline expansion "
    "integration certification training infrastructure deployment roadmap"
).split()

FOLDERS = ["Contracts", "Finance", "HR", "Operations", "Archive", "Misc"]

# Each case: a fact planted in exactly one file whose NAME gives nothing away.
# query_exact  — what a user types when they remember the wording
# query_fuzzy  — how a user asks when they only remember the concept
CASES = [
    dict(qid="Q01", folder="Misc", filename="scan_0231 copy.docx", kind="docx",
         fact="The aggregate liability cap under the Meridian master services agreement is INR 4.2 crore.",
         query_exact="liability cap Meridian master services agreement",
         query_fuzzy="maximum amount we could owe Meridian if something goes wrong",
         answer="INR 4.2 crore"),
    dict(qid="Q02", folder="Archive", filename="Final_v7_USE_THIS_ONE.xlsx", kind="xlsx",
         fact="FY27 travel budget approved ceiling 1,84,00,000 rupees for all business units combined.",
         query_exact="FY27 travel budget approved ceiling",
         query_fuzzy="how much can the whole company spend on travel next year",
         answer="1,84,00,000"),
    dict(qid="Q03", folder="Operations", filename="meeting notes (recovered).pdf", kind="pdf",
         fact="Data retention period for customer telemetry was fixed at 36 months by the security committee.",
         query_exact="data retention period customer telemetry",
         query_fuzzy="how long do we keep customer usage data before deleting it",
         answer="36 months"),
    dict(qid="Q04", folder="HR", filename="misc_old_2.docx", kind="docx",
         fact="Notice period for senior engineers was revised to 90 days effective January 2026.",
         query_exact="notice period senior engineers revised",
         query_fuzzy="how many days notice does a senior dev have to give before quitting",
         answer="90 days"),
    dict(qid="Q05", folder="Contracts", filename="Document1 (3).pdf", kind="pdf",
         fact="The Helios renewal includes a 7 percent uplift cap on annual subscription pricing.",
         query_exact="Helios renewal uplift cap subscription pricing",
         query_fuzzy="by how much can Helios raise our subscription price each year",
         answer="7 percent"),
    dict(qid="Q06", folder="Finance", filename="untitled spreadsheet (8).xlsx", kind="xlsx",
         fact="Procurement threshold requiring board approval is any purchase above 25 lakh rupees.",
         query_exact="procurement threshold board approval purchase",
         query_fuzzy="at what purchase size do we need the board to sign off",
         answer="25 lakh"),
    dict(qid="Q07", folder="Operations", filename="slides_backup_old.pptx", kind="pptx",
         fact="The Jaipur datacenter migration cutover is scheduled for the weekend of 14 September 2026.",
         query_exact="Jaipur datacenter migration cutover scheduled",
         query_fuzzy="when do we actually switch over to the new datacenter",
         answer="14 September 2026"),
    dict(qid="Q08", folder="Archive", filename="notes_final_FINAL.docx", kind="docx",
         fact="Warranty coverage for the Kestrel hardware line was extended to 60 months for enterprise tier.",
         query_exact="warranty coverage Kestrel hardware extended",
         query_fuzzy="how long is the warranty on Kestrel boxes for big customers",
         answer="60 months"),
    dict(qid="Q09", folder="HR", filename="export (1).pdf", kind="pdf",
         fact="Annual training stipend per employee is 75,000 rupees, claimable quarterly.",
         query_exact="annual training stipend per employee claimable",
         query_fuzzy="how much money does each person get for courses every year",
         answer="75,000"),
    dict(qid="Q10", folder="Contracts", filename="tmp_review_dont_delete.docx", kind="docx",
         fact="Payment terms with Northwind were renegotiated from net-30 to net-45 in March 2026.",
         query_exact="payment terms Northwind renegotiated net-45",
         query_fuzzy="how many days does Northwind have to pay us now",
         answer="net-45"),
    dict(qid="Q11", folder="Finance", filename="WhatsApp_attachment_004.pdf", kind="pdf",
         fact="The audit committee approved a contingency reserve of 2.1 percent of annual operating expenditure.",
         query_exact="audit committee contingency reserve operating expenditure",
         query_fuzzy="what slice of opex did we set aside for emergencies",
         answer="2.1 percent"),
    dict(qid="Q12", folder="Misc", filename="New folder summary v2.pptx", kind="pptx",
         fact="Escalation SLA for severity-1 incidents is 15 minutes to first response, around the clock.",
         query_exact="escalation SLA severity-1 incidents first response",
         query_fuzzy="how fast must someone respond when production is completely down",
         answer="15 minutes"),
]


def _sentences(rng: random.Random, n: int) -> str:
    out = []
    for _ in range(n):
        words = [rng.choice(WORDS) for _ in range(rng.randint(8, 16))]
        out.append((" ".join(words)).capitalize() + ".")
    return " ".join(out)


def _make_pdf(text: str) -> bytes:
    """One-page PDF with a Flate-compressed content stream — like real PDFs,
    so raw `grep` over the file cannot see the text."""
    safe = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    # Wrap at word boundaries (like real PDF layout) so words never split.
    lines, cur = [], ""
    for word in safe.split():
        if cur and len(cur) + 1 + len(word) > 90:
            lines.append(cur)
            cur = word
        else:
            cur = f"{cur} {word}" if cur else word
    if cur:
        lines.append(cur)
    lines = lines or [""]
    ops = ["BT /F1 10 Tf 40 760 Td 12 TL"]
    for ln in lines:
        ops.append(f"({ln}) Tj T*")
    ops.append("ET")
    stream = zlib.compress(" ".join(ops).encode("latin-1", errors="replace"))
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>"),
        (b"<< /Length " + str(len(stream)).encode() + b" /Filter /FlateDecode >>\n"
         b"stream\n" + stream + b"\nendstream"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objects, 1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n").encode()
    return bytes(out)


def _write_docx(path: Path, paragraphs) -> None:
    from docx import Document
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))


def _write_xlsx(path: Path, rows) -> None:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    for r in rows:
        ws.append(r)
    wb.save(str(path))


def _write_pptx(path: Path, slides) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        box.text_frame.text = text
    prs.save(str(path))


def _write_fact_file(path: Path, kind: str, fact: str, rng: random.Random) -> None:
    before = _sentences(rng, rng.randint(6, 12))
    after = _sentences(rng, rng.randint(6, 12))
    if kind == "docx":
        _write_docx(path, [before, fact, after])
    elif kind == "xlsx":
        _write_xlsx(path, [["topic", "detail"],
                           ["context", before[:180]],
                           ["decision", fact],
                           ["context", after[:180]]])
    elif kind == "pptx":
        _write_pptx(path, [before[:300], fact, after[:300]])
    elif kind == "pdf":
        path.write_bytes(_make_pdf(f"{before} {fact} {after}"))


def _write_distractor(path: Path, kind: str, rng: random.Random) -> None:
    body = _sentences(rng, rng.randint(8, 20))
    if kind == "docx":
        _write_docx(path, [body])
    elif kind == "xlsx":
        _write_xlsx(path, [["notes"], [body[:300]], [body[300:600]]])
    elif kind == "pptx":
        _write_pptx(path, [body[:300], body[300:600]])
    elif kind == "pdf":
        path.write_bytes(_make_pdf(body))
    else:
        path.write_text(body, encoding="utf-8")


def generate(root: Path, distractors: int = 150) -> dict:
    """Build the corpus. Returns {qid: case} with rel paths filled in."""
    rng = random.Random(SEED)
    root.mkdir(parents=True, exist_ok=True)
    for folder in FOLDERS:
        (root / folder).mkdir(exist_ok=True)

    cases = {}
    for case in CASES:
        rel = f"{case['folder']}/{case['filename']}"
        _write_fact_file(root / rel, case["kind"], case["fact"], rng)
        cases[case["qid"]] = dict(case, rel=rel)

    kinds = ["docx", "xlsx", "pptx", "pdf", "md", "txt", "csv"]
    names = ["report", "summary", "notes", "data", "export", "draft",
             "backup", "copy of copy", "v2 final", "untitled"]
    for i in range(distractors):
        kind = kinds[i % len(kinds)]
        folder = FOLDERS[i % len(FOLDERS)]
        name = f"{rng.choice(names)}_{i:03d}.{kind}"
        _write_distractor(root / folder / name, kind, rng)
    return cases


if __name__ == "__main__":
    import sys
    target = Path(sys.argv[1] if len(sys.argv) > 1 else "./bench_corpus")
    generated = generate(target)
    print(f"corpus at {target}: {len(generated)} planted facts, "
          f"{sum(1 for _ in target.rglob('*') if _.is_file())} files total")
