"""Format-aware text extraction.

Each extractor imports its dependency lazily so that a missing library only
affects that format. `.doc`/`.rtf` rely on macOS `textutil`; on other
platforms they are reported as unsupported rather than failing.
"""
from __future__ import annotations

import json
import logging
import os
import subprocess
import sys
import warnings
from pathlib import Path

# Quiet the noisy third-party extractor chatter: pdfminer emits hundreds of
# "Cannot set gray color … invalid float value" / "FontBBox" lines on real-world
# PDFs, and openpyxl warns about unsupported spreadsheet extensions. Neither is a
# failure (extraction proceeds), but the flood buries real errors. Set DOCDEX_DEBUG
# in the environment to restore the original verbosity.
if not os.environ.get("DOCDEX_DEBUG"):
    logging.getLogger("pdfminer").setLevel(logging.ERROR)

SECRETS_FILENAME = ".docdex.secrets.json"

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".html", ".htm", ".py", ".js", ".ts",
    ".tsx", ".css", ".yaml", ".yml", ".xml", ".log", ".bat", ".tex", ".rst",
}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".xlsm", ".pdf"}
TEXTUTIL_EXTENSIONS = {".doc", ".rtf"}  # macOS only

UNSUPPORTED_PREFIX = "[unsupported"

# What people type when they mean a file type. `ppt` is the important one: a user
# asking for "an atp ppt" means a deck, and the deck on disk is `.pptx`. Legacy
# `.ppt` is deliberately absent from EXTRACTABLE below — this codebase cannot read
# it — so mapping the word to `.pptx` is the difference between finding the deck
# and being told, truthfully but uselessly, that nothing matched.
EXT_ALIASES = {
    "ppt": ".pptx", "powerpoint": ".pptx", "deck": ".pptx", "slides": ".pptx",
    "doc": ".docx", "word": ".docx",
    "xls": ".xlsx", "excel": ".xlsx", "spreadsheet": ".xlsx",
    "markdown": ".md", "text": ".txt",
}


def extractable_extensions() -> set:
    """Every extension this build can actually read text out of."""
    return TEXT_EXTENSIONS | OFFICE_EXTENSIONS | TEXTUTIL_EXTENSIONS


def normalise_exts(values) -> tuple:
    """(extensions to filter on, values we could not honour).

    Returns the unusable ones rather than dropping them, because a filter that
    silently matches nothing is worse than no filter: it answers "absent" for a
    document that is present, and this tool is used on due-diligence questions
    where that is the expensive direction to be wrong in. The caller is expected
    to say so out loud.
    """
    known = extractable_extensions()
    wanted, unknown = set(), []
    for raw in values or []:
        token = str(raw).strip().lower().lstrip("*")
        if not token:
            continue
        token = EXT_ALIASES.get(token.lstrip("."), token)
        if not token.startswith("."):
            token = "." + token
        if token in known:
            wanted.add(token)
        else:
            unknown.append(str(raw))
    return wanted, unknown

# How many leading bytes to inspect when deciding whether a file is binary.
_SNIFF_BYTES = 8192
# C0 control bytes that are legitimate in real text: tab, LF, CR, form-feed.
_TEXT_CTRL = {0x09, 0x0a, 0x0d, 0x0c}
# Control chars stripped from extracted text: every C0 char that is not normal
# whitespace, plus DEL. Built once; used by sanitize_text via str.translate.
_STRIP_CTRL = {c: None for c in range(0x20) if c not in _TEXT_CTRL}
_STRIP_CTRL[0x7f] = None


def looks_binary(sample: bytes) -> bool:
    """Heuristic: does this byte sample look like binary (non-text) content?

    A NUL byte in the sniff window is the strong signal — it appears almost
    immediately in images, archives, executables, and random blobs, and never
    in real UTF-8 text. Failing that, a high proportion of C0 control bytes that
    never occur in text marks it binary. Bytes >= 0x80 are deliberately NOT
    counted, so heavily non-ASCII UTF-8 (Devanagari, accented Latin, currency
    symbols) is correctly treated as text rather than mistaken for binary.
    """
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    ctrl = sum(1 for b in sample if b < 0x20 and b not in _TEXT_CTRL)
    return ctrl / len(sample) > 0.30


def sanitize_text(s: str) -> str:
    """Strip NUL and stray control characters from extracted text while keeping
    tabs/newlines and all printable/Unicode content. Some PDF/office extractions
    leak the odd control char into otherwise-real text; removing them keeps the
    cache and index clean without discarding the document. Idempotent."""
    return s.translate(_STRIP_CTRL)


def textutil_available() -> bool:
    return sys.platform == "darwin"


def supported_extensions() -> set:
    exts = TEXT_EXTENSIONS | OFFICE_EXTENSIONS
    if textutil_available():
        exts |= TEXTUTIL_EXTENSIONS
    return exts


def is_supported(path) -> bool:
    return Path(path).suffix.lower() in supported_extensions()


def extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- SLIDE {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def extract_xlsx(path: str) -> str:
    import openpyxl
    parts = []
    with warnings.catch_warnings():
        if not os.environ.get("DOCDEX_DEBUG"):
            warnings.simplefilter("ignore")  # openpyxl "extension not supported" noise
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        try:
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"\n=== SHEET: {sheet} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
        finally:
            wb.close()
    return "\n".join(parts)


def read_secrets(root) -> dict:
    """Load the optional, user-controlled PDF-password map. Checks the v2 home
    location `<root>/.docdex/secrets.json` first, then the legacy
    `<root>/.docdex.secrets.json`. Missing or corrupt → empty dict (never
    raises). It lives inside the hidden home (or is a root dotfile), so the
    walker never indexes it; it is never committed to the docdex repo and its
    values are never logged."""
    root = Path(root)
    for candidate in (root / ".docdex" / "secrets.json", root / SECRETS_FILENAME):
        try:
            data = json.loads(candidate.read_text(encoding="utf-8"))
        except (OSError, ValueError):
            continue
        return data if isinstance(data, dict) else {}
    return {}


def candidate_passwords(rel_path: str, secrets: dict) -> list:
    """Passwords whose key is a substring of the file's path. An empty-string key
    matches every path (a deliberate corpus-wide fallback)."""
    return [pw for key, pw in secrets.items() if key in rel_path]


# The OLE2 / "CDFV2" container. A password-protected .docx/.pptx/.xlsx is written as
# one of these rather than as a zip, so it raises the same "not a package" error as a
# file that is simply damaged — the two need telling apart before either can be acted on.
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _head_bytes(path: str, count: int) -> bytes:
    try:
        with open(path, "rb") as f:
            return f.read(count)
    except OSError:
        return b""


def describe_failure(exc: BaseException, path: str, passwords=()) -> str:
    """Plain-language reason an extraction failed, for `extract_status` and `doctor`.

    The raw exception text was actively misleading on the real corpus. A present but
    truncated `.docx` reported `PackageNotFoundError: Package not found at '<path>'`,
    which reads as "that file is missing" and sends you looking for a path problem; an
    encrypted PDF with no configured password reported `PDFPasswordIncorrect:` with an
    empty message, which reads as "docdex tried a password and got it wrong". Telling
    six damaged files from four protected ones took a by-hand investigation of zip
    directories and PDF trailers that this function now does up front.

    The exception name is kept in brackets at the end so nothing is lost for debugging.
    Note that `passwords` only ever apply to PDFs (see `extract`), so an encrypted
    Office file must NOT be told to add one — advice that cannot work is the thing
    this is fixing.
    """
    name = type(exc).__name__
    if name == "PDFPasswordIncorrect":
        if passwords:
            return (f"encrypted PDF — the password configured in .docdex/secrets.json "
                    f"did not work [{name}]")
        return (f"encrypted PDF — no password configured; add one to "
                f".docdex/secrets.json [{name}]")
    # `BadZipFile` belongs with `PackageNotFoundError`: python-docx and python-pptx
    # wrap a non-zip file in their own error, but openpyxl lets zipfile's through, so
    # an encrypted or truncated .xlsx used to fall past every branch here and reach the
    # user as a bare "BadZipFile: File is not a zip file". Found by adversarial review
    # of this release's tests, which only ever exercised .docx.
    if name in ("PackageNotFoundError", "BadZipFile", "PSEOF") \
            and not os.path.exists(path):
        # Adversarial review of this release: both branches below assert "the file is
        # present", which is the whole point of them — but a file deleted between the
        # inventory walk and extraction raises the same errors, and then the message
        # states the exact opposite of the filesystem. Checked, not assumed.
        return f"file disappeared during extraction — it is no longer on disk [{name}]"
    if name in ("PackageNotFoundError", "BadZipFile"):
        if _head_bytes(path, len(_OLE2_MAGIC)).startswith(_OLE2_MAGIC):
            # Two different causes land here and the remedy is the same for both, so
            # the container is named rather than the cause guessed: a password-protected
            # Office file is written as OLE2/CDFV2, and so is a legacy binary document
            # that has been given a modern extension. Adversarial review flagged this
            # branch for asserting "encrypted"; measurement showed a genuine .doc/.xls
            # never reaches it (those route to textutil and extract normally).
            #
            # `secrets.json` is deliberately NOT mentioned. It used to appear here as a
            # disclaimer ("passwords work for PDFs only"), and review pointed out that
            # naming a file docdex cannot use for this format is an invitation to try
            # it whatever the surrounding words say.
            return (f"encrypted or legacy Office file — an OLE2/CDFV2 container, not a "
                    f".docx/.pptx/.xlsx package; docdex cannot read password-protected "
                    f"Office files at all — re-save it as an unprotected "
                    f".docx/.pptx/.xlsx [{name}]")
        return (f"damaged or truncated Office file — the file is present but has no "
                f"complete zip directory [{name}]")
    if name == "PSEOF":
        return (f"damaged or truncated PDF — the file is present but ends mid-structure "
                f"[{name}]")
    return f"{name}: {exc}"


def extract_pdf(path: str, passwords=()) -> str:
    from pdfminer.high_level import extract_text
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    last_err = None
    for pw in ("", *passwords):  # try unencrypted / owner-readable first, then candidates
        try:
            return extract_text(path, password=pw) or ""
        except PDFPasswordIncorrect as e:
            last_err = e
    raise last_err  # encrypted and no candidate password worked


def extract_plain(path: str) -> str:
    """Read a whitelisted text file, but refuse binary content.

    A file can carry a text extension (.log/.csv/.json/.xml/.txt/.html) yet hold
    binary data — a log rotated into a blob, a renamed image, a database dump.
    Reading that with errors="replace" produced megabytes of replacement-char
    garbage that got chunked and embedded; at size/scale it also hung sync. We
    sniff the first block and report binary content as unsupported instead."""
    with open(path, "rb") as f:
        head = f.read(_SNIFF_BYTES)
        if looks_binary(head):
            return f"{UNSUPPORTED_PREFIX} binary content ({Path(path).suffix.lower()})]"
        rest = f.read()
    return (head + rest).decode("utf-8", errors="replace")


def extract_with_textutil(path: str) -> str:
    r = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", path],
        capture_output=True, text=True,
    )
    if r.returncode != 0:
        raise RuntimeError(r.stderr.strip() or f"textutil exited {r.returncode}")
    return r.stdout


def extract(path, passwords=()) -> str:
    """Return extracted text, or an `[unsupported ...]` marker string.

    `passwords` are candidate passwords tried (in order) for an encrypted PDF;
    callers build them with `read_secrets` + `candidate_passwords`."""
    p = Path(path)
    ext = p.suffix.lower()
    if ext == ".docx":
        return extract_docx(str(p))
    if ext == ".pptx":
        return extract_pptx(str(p))
    if ext in (".xlsx", ".xlsm"):
        return extract_xlsx(str(p))
    if ext == ".pdf":
        return extract_pdf(str(p), passwords=passwords)
    if ext in TEXTUTIL_EXTENSIONS:
        if textutil_available():
            return extract_with_textutil(str(p))
        return f"{UNSUPPORTED_PREFIX} {ext} requires macOS textutil; convert to .docx]"
    if ext in TEXT_EXTENSIONS:
        return extract_plain(str(p))
    return f"{UNSUPPORTED_PREFIX} extension: {ext}]"
